from docx import Document
from pptx import Presentation


def transliterate(text):
    diction = {'А': 'A', 'Б': 'B', 'В': 'V', 'Г': 'G', 'Д': 'D', 'Е': 'E', 'Ё': 'E',
               'Ж': 'Zh', 'З': 'Z', 'И': 'I', 'Й': 'I', 'К': 'K', 'Л': 'L', 'М': 'M',
               'Н': 'N', 'О': 'O', 'П': 'P', 'Р': 'R', 'С': 'S', 'Т': 'T', 'У': 'U',
               'Ф': 'F', 'Х': 'Kh', 'Ц': 'Tc', 'Ч': 'Ch', 'Ш': 'Sh', 'Щ': 'Shch',
               'Ы': 'Y', 'Э': 'E', 'Ю': 'Iu', 'Я': 'Ia', 'а': 'a', 'б': 'b', 'в': 'v',
               'г': 'g', 'д': 'd', 'е': 'e', 'ё': 'e', 'ж': 'zh', 'з': 'z', 'и': 'i',
               'й': 'i', 'к': 'k', 'л': 'l', 'м': 'm', 'н': 'n', 'о': 'o', 'п': 'p',
               'р': 'r', 'с': 's', 'т': 't', 'у': 'u', 'ф': 'f', 'х': 'kh', 'ц': 'tc',
               'ч': 'ch', 'ш': 'sh', 'щ': 'shch', 'ы': 'y', 'э': 'e', 'ю': 'iu', 'я': 'ia'
               }
    res = ''
    for i in text:
        try:
            if i == 'ъ' or i == 'ь' or i == 'Ъ' or i == 'Ь':
                continue
            else:
                res += diction[i]
        except KeyError:
            res += i

    return res


class TicketSystem:
    def __init__(self):
        self.sessions = []
        self.cinemas = []
        self.cinema_halls = []
        self.info = {}

    def check_cinema(self, z):
        for cinema in self.cinemas:
            if z == cinema.name:
                return cinema

    def check_hall(self, c1):
        for cinema in self.cinemas:
            if c1 == cinema.name:
                for hall in cinema:
                    if hall.name == hall:
                        return hall
        return None

    def add_cinema(self, other):
        self.cinemas.append(other)

    def add_hall(self, hall, cinema):
        for cin_t in range(len(self.cinemas)):
            if self.cinemas[cin_t] == cinema:
                self.cinemas[cin_t].halls.append(hall)
                self.info[cinema] = {hall: ''}
                print(1)
                return
        print('Ошибка! Такого кинотеатра нет в системе.')

    def add_session(self, s, cinema, hall):
        for cin_t in range(len(self.cinemas)):
            if self.cinemas[cin_t] == cinema:
                for hl_t in range(len(self.cinemas[cin_t].halls)):
                    if self.cinemas[cin_t].halls[hl_t] == hall:
                        if s not in hall.sessions:
                            self.cinemas[cin_t].halls[hl_t].sessions.append(s)
                            return
                print('Ошибка! Такого зала нет в кинотеатре.')
                return
        print('Ошибка! Такого кинотеатра нет в системе.')

    def __str__(self):
        information = ''
        for cinema in self.cinemas:
            information += (cinema.name + ':\n')
            for hall in cinema.halls:
                information += f'\t{hall.name}:\n'
                for session in hall.sessions:
                    information += f'\t\t{str(session)}\n'
        return information


class Session:
    def __init__(self, name, time, price, cinema, hall):
        self.name, self.time, self.price = name, time, price
        self.not_empty_places = []
        self.hall = hall.name
        self.cinema = cinema.name
        self.x, self.y = hall.length, hall.width
        self.empty_places = self.x * self.y

    def is_empty(self, x, y):
        if (x, y) not in self.not_empty_places and \
                (0 < x <= self.x) and (0 < y <= self.y):
            return True
        else:
            return False

    def buy_ticket(self, x, y):
        if self.is_empty(x, y):
            print('Билет куплен.')
            self.not_empty_places.append((x, y))
            self.empty_places -= 1
        else:
            print('Место уже занято или недействительно.')

    def __str__(self):
        return f'Фильм: {self.name}, Время: {self.time}, Цена: {self.price}, ' \
               f'Свободных мест: {self.empty_places}'

    def full_information(self):
        return [f'Кинотеатр: {self.cinema}', f'Зал: {self.hall}', f'Фильм: {self.name}', f'Время: {self.time}',
                f'Цена: {self.price}', f'Свободных мест: {self.empty_places}']

    def create_commercial_flyer(self):
        name = transliterate('_'.join(self.name.split())) + '-' + transliterate('_'.join(self.cinema.split())) + '-' + \
               transliterate('_'.join(self.hall.split())) + '-' + ''.join(self.time.split(':')) + '.docx'
        commercial = Document()
        info = self.full_information()
        commercial.add_heading(self.name, 0)
        commercial.add_paragraph(f'Администрация кинотеатра "Мир" приглашает вас на сеанс фильма '
                                 f'{self.name}! Получите максимальное количество восторга и обескураженности всего лишь'
                                 f' за {self.price} рублей! Показ сеанса будет проходить в {self.hall}в '
                                 f'{self.time}. Поторопись, ведь сталось всего {self.empty_places} мест!\nЖдём тебя на '
                                 f'показе!\n С уважением, администрация кинотеатра "{self.cinema}"')
        commercial.add_heading('Краткая информация', level=1)
        commercial.add_paragraph(info[0],
                                 style='List Bullet').bold = True
        commercial.add_paragraph(info[1],
                                 style='List Bullet').bold = True
        commercial.add_paragraph(info[2],
                                 style='List Bullet').bold = True
        commercial.add_paragraph(info[3],
                                 style='List Bullet').bold = True
        commercial.add_paragraph(info[4],
                                 style='List Bullet').bold = True
        commercial.add_paragraph(info[5],
                                 style='List Bullet').bold = True
        commercial.save(name)
        print('Листовка успешно создана!')


class Hall:  # Класс, отвечающий за залы просмотра
    def __init__(self, name, length, width):
        self.name, self.length, self.width = name, length, width
        self.sessions = []

    def check_session(self, n_t_p):
        for session in self.sessions:
            if session.name + ' ' + session.time + ' ' + session.price == n_t_p:
                return session
        return None


class Cinema:
    def __init__(self, name):
        self.name = name
        self.halls = []

    def check_hall(self, name):
        for hall in self.halls:
            if hall.name == name:
                return hall
        return None

    def make_presentation(self):
        prs = Presentation()
        information = ''
        for hall in self.halls:
            title_slide_layout = prs.slide_layouts[2]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = hall.name
            for session in hall.sessions:
                information += f'\n{session.name}\n{str(session)}'
            if not len(information):
                information = 'В зале нет сеансов'
            subtitle.text = information
            information = ''
        prs.save(f'{transliterate(self.name)}.pptx')
