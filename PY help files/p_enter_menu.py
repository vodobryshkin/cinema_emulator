import sys
from docx import Document
from pptx import Presentation
from datetime import date as dt
from datetime import datetime
from random import choice
import sqlite3
from PyQt5.QtWidgets import QApplication, QMainWindow, QWidget, QTableWidgetItem, QFileDialog
from PyQt5.QtGui import QPixmap
from ui_authorization import Ui_AuthorisationForm
from ui_start import Ui_StartWindow
from ui_registration import Ui_RegistrationWindow
from ui_admin_profile import Ui_AdminModeWindow
from ui_create_hall import Ui_CreateHallWindow
from ui_create_session import Ui_CreateSessionWindow
from ui_create_product import Ui_CreateProductWindow
from ui_create_commercial import Ui_CreateCommercialWindow
from ui_user_profile import Ui_UserModeWindow


class CurrentHallError(Exception):
    pass


class CurrentTimeError(Exception):
    pass


sym = '1 2 3 4 5 6 7 8 9 0 q w e r t y u i o p a s d f g h j k l z x c v b n m ! @ # $ % ^ & * ( )'.split()

base = sqlite3.connect('db_cinema.db')
cur = base.cursor()

base.execute(
    'CREATE TABLE IF NOT EXISTS users (id INTEGER PRIMARY KEY AUTOINCREMENT NOT NULL, login TEXT NOT NULL UNIQUE'
    ', password TEXT NOT NULL)')

try:
    cur.execute('INSERT INTO users VALUES(?, ?, ?)', (None, 'admin', '1'))
except sqlite3.IntegrityError:
    pass

base.execute(
    'CREATE TABLE IF NOT EXISTS halls (name TEXT NOT NULL UNIQUE, size TEXT NOT NULL, employment INTEGER,'
    ' id INTEGER PRIMARY KEY AUTOINCREMENT NOT NULL)')

base.execute(
    'CREATE TABLE IF NOT EXISTS sessions (title TEXT, hall TEXT, date TEXT, '
    'price INTEGER, time TEXT, empty_places INTEGER, id  INTEGER NOT NULL PRIMARY KEY AUTOINCREMENT)')

base.execute('CREATE TABLE IF NOT EXISTS market (title TEXT UNIQUE, price INTEGER, amount INTEGER, '
             'id INTEGER PRIMARY KEY AUTOINCREMENT)')

base.commit()


def transliterate(text):
    diction = {'А': 'A', 'Б': 'B', 'В': 'V', 'Г': 'G', 'Д': 'D', 'Е': 'E', 'Ё': 'E',
               'Ж': 'Zh', 'З': 'Z', 'И': 'I', 'Й': 'I', 'К': 'K', 'Л': 'L', 'М': 'M',
               'Н': 'N', 'О': 'O', 'П': 'P', 'Р': 'R', 'С': 'S', 'Т': 'T', 'У': 'U',
               'Ф': 'F', 'Х': 'Kh', 'Ц': 'Tc', 'Ч': 'Ch', 'Ш': 'Sh', 'Щ': 'Shch',
               'Ы': 'Y', 'Э': 'E', 'Ю': 'Iu', 'Я': 'Ia', 'а': 'a', 'б': 'b', 'в': 'v',
               'г': 'g', 'д': 'd', 'е': 'e', 'ё': 'e', 'ж': 'zh', 'з': 'z', 'и': 'i',
               'й': 'i', 'к': 'k', 'л': 'l', 'м': 'm', 'н': 'n', 'о': 'o', 'п': 'p',
               'р': 'r', 'с': 's', 'т': 't', 'у': 'u', 'ф': 'f', 'х': 'kh', 'ц': 'tc',
               'ч': 'ch', 'ш': 'sh', 'щ': 'shch', 'ы': 'y', 'э': 'e', 'ю': 'iu', 'я': 'ia'
               }
    res = ''
    for i in text:
        try:
            if i == 'ъ' or i == 'ь' or i == 'Ъ' or i == 'Ь':
                continue
            else:
                res += diction[i]
        except KeyError:
            res += i

    return res


def data_check(data):
    try:
        data = data.split('.')
    except TypeError:
        return False

    try:
        int(data[0]), int(data[1]), int(data[2])
    except ValueError:
        return False

    d, m, y = data[0], data[1], data[2]
    return len(d) == len(m) == 2 and len(y) == 4


def time_check(t):
    try:
        t = t.split(':')
    except TypeError:
        return False
    try:
        int(t[0]), int(t[1])
    except ValueError:
        return False

    h, m = t[0], t[1]
    return len(h) == len(m) == 2


class StartWindow(QWidget, Ui_StartWindow):
    def __init__(self):
        super().__init__()
        self.setupUi(self)

        self.pixmap = QPixmap('mir_logo.jpg')
        self.photo_label.setPixmap(self.pixmap)

        self.log_in_widget = AuthorizationWindow(self)
        self.log_in_widget.back_button.clicked.connect(self.log_in_back)

        self.registration_widget = RegistrationWindow(self)
        self.registration_widget.back_button.clicked.connect(self.registration_back)

        self.log_in_button.clicked.connect(self.to_log_in)
        self.registration_button.clicked.connect(self.to_registration)

    def to_log_in(self):
        self.log_in_widget.show()
        self.hide()

    def to_registration(self):
        self.registration_widget.show()
        self.hide()

    def log_in_back(self):
        self.log_in_widget.hide()
        self.show()

    def registration_back(self):
        self.registration_widget.hide()
        self.show()


class AuthorizationWindow(QWidget, Ui_AuthorisationForm):
    def __init__(self, parent):
        super().__init__()
        self.setupUi(self)
        self.back = False
        self.parent = parent
        self.admin_window = AdminModeWindow(self)
        self.user_window = UserModeWindow(self)

        self.confirm_button.clicked.connect(self.check_user)

    def check_user(self):
        login = self.enter_login.text()
        password = self.enter_password.text()
        cur.execute(f"SELECT login from users WHERE login='{login}' AND password='{password}';")
        if not cur.fetchone():
            self.wrong_label.setText('Неверный логин или пароль.')
        else:
            self.hide()
            if self.enter_login.text() == 'admin':
                self.admin_window.show()
            else:
                self.user_window.show()
                self.user_window.welcome_label.setText(f'Добро пожаловать, {login}.')

    def go_back(self):
        self.hide()


class RegistrationWindow(QWidget, Ui_RegistrationWindow):
    def __init__(self, parent):
        super().__init__()
        self.setupUi(self)
        self.back = False
        self.parent = parent

        self.confirm_button.clicked.connect(self.check_user)

    def check_user(self):
        login = self.enter_login.text()
        password = self.enter_password.text()
        try:
            cur.execute('INSERT INTO users VALUES(?, ?, ?)', (None, login, password))
            self.output_label.setText('Регистрация прошла успешно.')
        except sqlite3.IntegrityError:
            self.output_label.setText('Пользователь с таким логином\nуже существует.')

        base.commit()


class CreateHallWindow(QWidget, Ui_CreateHallWindow):
    def __init__(self, parent):
        super().__init__()
        self.setupUi(self)
        self.parent = parent

        self.create_hall_button.clicked.connect(self.create_hall)

    def create_hall(self):
        create_complete = False
        size = f"'{str(self.width_input.text()) + '*' + str(self.length_input.text())}'"
        try:
            cur.execute("""INSERT INTO halls VALUES (?, ?, ?, ?)""", (self.name_input.text(), size, 0, None))
            self.error_widget.setText('Зал успешно добавлен.')
            create_complete = True
            base.commit()
            self.hide()
        except sqlite3.IntegrityError:
            self.error_widget.setText('Зал с таким названием уже есть.')

        if create_complete:
            self.parent.hall_table_update()


class CreateSessionWindow(QWidget, Ui_CreateSessionWindow):
    def __init__(self, parent):
        super().__init__()
        self.setupUi(self)
        self.parent = parent
        self.chosen = ''

        self.create_session_button.clicked.connect(self.create_session)
        self.combo_box.activated.connect(self.handle_activated)

    def handle_activated(self, index):
        self.chosen = self.combo_box.itemText(index)

    def create_session(self):
        name = self.name_input.text()
        hall = self.chosen
        k = cur.execute(f"""SELECT size FROM halls WHERE name='{self.chosen}'""").fetchone()[0]
        data = self.date_edit.dateTime().toString("dd.MM.yyyy")
        t = self.time_edit.time().toString()
        k = k.split('*')
        price = self.spin_box.value()

        n = cur.execute(f"""SELECT employment FROM halls WHERE name='{self.chosen}'""").fetchone()[0] + 1
        cur.execute("""INSERT INTO sessions VALUES (?, ?, ?, ?, ?, ?, ?)""", (name, hall, data, price, t,
                                                                              int(k[0][1:]) * int(k[1][:-1]), None))
        cur.execute(f"""UPDATE halls SET employment = {n} WHERE name='{self.chosen}'""")
        self.name_input.text()
        self.error_widget.setText('Сеанс успешно добавлен.')
        base.commit()
        self.hide()

        self.parent.hall_table_update()
        self.parent.session_table_update()


class CreateProductWindow(QWidget, Ui_CreateProductWindow):
    def __init__(self, parent):
        super().__init__()
        self.setupUi(self)
        self.parent = parent

        self.create_product_button.clicked.connect(self.create_product)

    def create_product(self):
        create_complete = False
        try:
            price = int(self.price_input.text())
            amount = int(self.anum_input.text())
            if not cur.execute(f"""SELECT * FROM market WHERE title='{self.name_input.text()}'""").fetchall():
                cur.execute("""INSERT INTO market VALUES (?, ?, ?, ?)""", (self.name_input.text(), price, amount, None))
                self.error_widget.setText('Товар успешно добавлен.')
                create_complete = True
                base.commit()
                self.hide()
            else:
                create_complete = True
                cur.execute(f"""UPDATE market SET price='{price}' WHERE title='{self.name_input.text()}'""")
                cur.execute(f"""UPDATE market SET amount='{amount}' WHERE title='{self.name_input.text()}'""")
                base.commit()
                self.hide()
        except ValueError:
            self.error_widget.setText('Неверный тип данных.')

        if create_complete:
            self.parent.market_table_update()


class CreateCommercialWindow(QWidget, Ui_CreateCommercialWindow):
    def __init__(self, parent):
        super().__init__()
        self.setupUi(self)
        self.parent = parent
        self.chosen = ''

        self.combo_box.activated.connect(self.handle_activated)
        self.create_button.clicked.connect(self.create_commercial)

    def handle_activated(self, index):
        self.chosen = self.combo_box.itemText(index)

    def create_commercial(self):
        data = cur.execute(f"""SELECT * FROM sessions WHERE title LIKE '{self.chosen.split('id:')[0]}' AND 
id={self.chosen.split('id:')[1]}""").fetchone()
        name = data[0]
        hall = cur.execute(f"""SELECT name FROM halls WHERE id={data[1]}""").fetchone()[0]
        date = data[2]
        price = str(data[3])
        t = data[4]
        empty_places = data[5]
        info = [f'Кинотеатр: "Мир"', f'Зал: {hall}', f'Фильм: {name}', f'Дата: {date}', f'Время: {t}',
                f'Цена: {price}', f'Свободных мест: {empty_places}']
        file_name = transliterate('_'.join(name.split())) + '-' + transliterate('_'.join('Мир'.split())) + '-'
        file_name += transliterate('_'.join(hall.split())) + '-' + transliterate('_'.join(date.split('.')))
        file_name += ''.join(t.split(':')) + '.docx'
        commercial = Document()
        commercial.add_heading(name, 0)
        commercial.add_paragraph(f'Администрация кинотеатра "Мир" приглашает вас на сеанс фильма '
                                 f'{name}! Получите максимальное количество восторга и обескураженности всего лишь'
                                 f' за {price} рублей! Показ сеанса будет проходить в {hall} {date} в'
                                 f'{t}. Поторопись, ведь сталось всего {empty_places} мест!\nЖдём тебя на '
                                 f'показе!\n С уважением, администрация кинотеатра "Мир"')
        commercial.add_heading('Краткая информация', level=1)
        commercial.add_paragraph(info[0],
                                 style='List Bullet').bold = True
        commercial.add_paragraph(info[1],
                                 style='List Bullet').bold = True
        commercial.add_paragraph(info[2],
                                 style='List Bullet').bold = True
        commercial.add_paragraph(info[3],
                                 style='List Bullet').bold = True
        commercial.add_paragraph(info[4],
                                 style='List Bullet').bold = True
        commercial.add_paragraph(info[5],
                                 style='List Bullet').bold = True
        commercial.add_paragraph(info[6],
                                 style='List Bullet').bold = True

        way = QFileDialog.getExistingDirectory(self, "Выберите путь", "C:\\Users\\")
        commercial.save(way + '/' + file_name)


class AdminModeWindow(QMainWindow, Ui_AdminModeWindow):
    def __init__(self, parent):
        super().__init__()
        self.setupUi(self)
        self.parent = parent

        self.create_hall_widget = CreateHallWindow(self)
        self.create_session_widget = CreateSessionWindow(self)
        self.create_product_widget = CreateProductWindow(self)
        self.create_commercial_widget = CreateCommercialWindow(self)

        self.hall_make_button.clicked.connect(self.create_hall)
        self.session_make_button.clicked.connect(self.create_session)
        self.add_product_button.clicked.connect(self.create_product)
        self.create_commercial_button.clicked.connect(self.create_commercial)
        self.create_feedback_button.clicked.connect(self.create_info)
        self.delete_session_button.clicked.connect(self.delete_session)
        self.back_button.clicked.connect(self.get_back)

        self.hall_table_update()
        self.session_table_update()
        self.market_table_update()

    def get_back(self):
        self.hide()
        self.parent.show()

    def delete_session(self):
        cur.execute(f"""DELETE FROM sessions WHERE id={self.spin_box.value()}""")
        base.commit()
        self.session_table_update()

    def market_table_update(self):
        query = "SELECT * FROM market"
        res = cur.execute(query).fetchall()

        self.market_table.setColumnCount(4)
        self.market_table.setHorizontalHeaderLabels(['Название', 'Цена', 'Кол-во', 'ID'])
        self.market_table.setRowCount(0)

        for i, row in enumerate(res):

            self.market_table.setRowCount(
                self.market_table.rowCount() + 1)

            for j, elem in enumerate(row):
                self.market_table.setItem(
                    i, j, QTableWidgetItem(str(elem)))

    def hall_table_update(self):
        query = "SELECT * FROM halls"
        res = cur.execute(query).fetchall()

        self.halls_table.setColumnCount(4)
        self.halls_table.setHorizontalHeaderLabels(['Название', 'Размер', 'Занятость', 'ID'])
        self.halls_table.setRowCount(0)

        for i, row in enumerate(res):
            self.halls_table.setRowCount(
                self.halls_table.rowCount() + 1)
            for j, elem in enumerate(row):
                self.halls_table.setItem(
                    i, j, QTableWidgetItem(str(elem)))

    def session_table_update(self):
        query = "SELECT * FROM sessions"
        res = cur.execute(query).fetchall()

        self.sessions_table.setColumnCount(7)
        self.sessions_table.setHorizontalHeaderLabels(['Название', 'Зал', 'Дата', 'Цена', 'Время',
                                                       '       Свободные места', 'ID'])
        self.sessions_table.setRowCount(0)

        for i, row in enumerate(res):
            self.sessions_table.setRowCount(
                self.sessions_table.rowCount() + 1)
            for j, elem in enumerate(row):
                self.sessions_table.setItem(
                    i, j, QTableWidgetItem(str(elem)))

    def create_hall(self):
        self.create_hall_widget.show()

    def create_session(self):
        self.create_session_widget.show()
        f = True
        self.create_session_widget.combo_box.clear()
        result = cur.execute("""SELECT name FROM halls""").fetchall()
        for i in result:
            if f:
                self.create_session_widget.chosen = f'{i[0]}'
                f = False
            self.create_session_widget.combo_box.addItem(f'{i[0]}')

    def create_product(self):
        self.create_product_widget.show()

    def create_commercial(self):
        self.create_commercial_widget.show()
        f = True
        self.create_commercial_widget.combo_box.clear()
        result = cur.execute("""SELECT title, id FROM sessions""").fetchall()
        for i in result:
            if f:
                self.create_commercial_widget.chosen = f'{i[0]}id:{i[1]}'
                f = False
            self.create_commercial_widget.combo_box.addItem(f'{i[0]}id:{i[1]}')

    def create_info(self):
        prs = Presentation()
        information = ''
        halls = [i[0] for i in cur.execute("""SELECT name FROM halls""").fetchall()]
        print(halls)
        for hall in halls:

            title_slide_layout = prs.slide_layouts[2]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = hall
            sessions = cur.execute(f"""SELECT title, id FROM sessions WHERE hall LIKE '{hall}'""").fetchall()

            for session in sessions:
                data = cur.execute(f"""SELECT * FROM sessions WHERE title LIKE '{session[0]}' AND 
                id={session[1]}""").fetchone()
                name = data[0]
                date = data[2]
                price = str(data[3])
                t = data[4]
                empty_places = data[5]
                info = f'Кинотеатр: "Мир", Зал: {hall}, Фильм: {name}, Дата: {date}', \
                       f'Время: {t}, Цена: {price}, Свободных мест: {empty_places}'
                information += f'\n{session}\n{info}'
            if not len(information):
                information = 'В зале нет сеансов'
            subtitle.text = information
            information = ''

        now = datetime.now()
        c_time = str(dt.today()) + "-".join(str(now.time()).split(':'))

        file_name = f"{transliterate('Обзор-' + c_time)}.pptx"

        way = str(QFileDialog.getExistingDirectory(self, "Выберите путь", "C:\\Users\\"))
        prs.save(way + '/' + file_name)


class UserModeWindow(QMainWindow, Ui_UserModeWindow):
    def __init__(self, parent):
        super().__init__()
        self.setupUi(self)

        self.parent = parent

        self.pixmap = QPixmap('cinema.jpg')
        self.photo_label.setPixmap(self.pixmap)

        self.chosen1 = ''
        self.chosen2 = ''

        self.market_table_update()
        self.session_table_update()

        self.combo_box_1.activated.connect(self.handle_activated_1)
        self.combo_box_2.activated.connect(self.handle_activated_2)

        self.update_button.clicked.connect(self.market_table_update)
        self.update_button.clicked.connect(self.session_table_update)
        self.buy_ticket_button.clicked.connect(self.buy_ticket)
        self.buy_product_button.clicked.connect(self.buy_product)
        self.back_button.clicked.connect(self.get_back)

    def get_back(self):
        self.hide()
        self.parent.show()

    def buy_product(self):
        res = cur.execute(f"""SELECT amount, price FROM market WHERE title='{self.chosen2.split()[0]}'""").fetchone()
        prev_a, price = res[0], res[1]
        try:
            if prev_a - int(self.input_val.text()) >= 0:
                cur.execute(f"""UPDATE market SET amount={prev_a - int(self.input_val.text())}
                                 WHERE title='{self.chosen2.split()[0]}'""")
                self.market_table_update()

                special_number = ''
                for _ in range(16):
                    special_number += choice(sym)

                now = datetime.now()
                way = str(QFileDialog.getExistingDirectory(self, "Выберите путь для сохранения чека", "C:\\Users\\"))

                f = open(way + '/' + 'ЧЕК от ' + str(now.date()) + '-'.join(str(now.time()).split(':'))
                         + '.txt', encoding='utf-8', mode='w')
                f.write(f'Благодарим за покупку предметов в магазине нашего кинотеатра.\n'
                        f'ИНФОРМАЦИЯ О ПЛАТЕЖЕ: дата - {now.date()}, время - {now.time()}.\n'
                        f'Сумма сделки - {int(self.input_val.text()) * price} рубля(ей).\n'
                        f'Для получения товара в кинотеатре, покажите кассиру специальный код: {special_number}.\n'
                        f'С уважением, администрация кинотеатра.')
                f.close()

                self.market_wrong_label.setText(f'Сделка на сумму {int(self.input_val.text()) * price} рубля(ей)'
                                                f' совершена.')
                base.commit()
            else:
                self.market_wrong_label.setText('Склад не может позволить себе эту операцию.')
        except ValueError:
            self.market_wrong_label.setText('Вводите цифры.')

    def buy_ticket(self):
        print(self.chosen1.split('id:')[0], self.chosen1.split('id:')[1])
        res = cur.execute(f"""SELECT empty_places, price FROM sessions WHERE title='{self.chosen1.split('id:')[0]}'
         and id={self.chosen1.split('id:')[1]}""").fetchone()

        print(res)

        prev_e, price = res[0], res[1]

        try:
            if prev_e - int(self.input_tick.text()) >= 0:
                cur.execute(f"""UPDATE sessions SET empty_places={prev_e - int(self.input_tick.text())}
                         WHERE title='{self.chosen1.split('id:')[0]}'""")
                print(1)
                self.session_table_update()

                special_number = ''
                for _ in range(16):
                    special_number += choice(sym)

                now = datetime.now()
                print(2)
                way = str(QFileDialog.getExistingDirectory(self, "Выберите путь для сохранения чека", "C:\\Users\\"))

                f = open(way + '/' + 'ЧЕК от ' + str(now.date()) + ' ' + '-'.join(str(now.time()).split(':'))
                         + '.txt', encoding='utf-8', mode='w')
                f.write(f'Благодарим за покупку билетов в нашем кинотеатре.\n'
                        f'ИНФОРМАЦИЯ О ПЛАТЕЖЕ: дата - {now.date()}, время - {now.time()}.\n'
                        f'Сумма сделки - {int(self.input_tick.text()) * price} рубля(ей).\n'
                        f'Для получения товара в кинотеатре, покажите кассиру специальный код: {special_number}.\n'
                        f'С уважением, администрация кинотеатра.')
                f.close()

                self.ticket_wrong_label.setText(f'Сделка на сумму {int(self.input_tick.text()) * price} рубля(ей)'
                                                f' совершена.')
                base.commit()
            else:
                self.ticket_wrong_label.setText('Касса не может позволить себе эту операцию.')

        except ValueError:
            self.ticket_wrong_label.setText('Вводите цифры.')

    def market_table_update(self):
        query = "SELECT * FROM market"
        res = cur.execute(query).fetchall()

        self.market_table.setColumnCount(4)
        self.market_table.setHorizontalHeaderLabels(['Название', 'Цена', 'Кол-во', 'ID'])
        self.market_table.setRowCount(0)

        for i, row in enumerate(res):

            self.market_table.setRowCount(
                self.market_table.rowCount() + 1)

            for j, elem in enumerate(row):
                self.market_table.setItem(
                    i, j, QTableWidgetItem(str(elem)))

        f = True
        self.combo_box_2.clear()
        result = cur.execute("""SELECT title, price FROM market""").fetchall()
        for i in result:
            if f:
                self.chosen2 = f'{i[0]} цена:{i[1]}рублей'
                f = False
            self.combo_box_2.addItem(f'{i[0]} цена:{i[1]}рублей')

    def session_table_update(self):
        query = "SELECT * FROM sessions"
        res = cur.execute(query).fetchall()

        self.sessions_table.setColumnCount(7)
        self.sessions_table.setHorizontalHeaderLabels(['Название', 'Зал', 'Дата', 'Цена', 'Время',
                                                       '       Свободные места', 'ID'])
        self.sessions_table.setRowCount(0)

        for i, row in enumerate(res):
            self.sessions_table.setRowCount(
                self.sessions_table.rowCount() + 1)
            for j, elem in enumerate(row):
                self.sessions_table.setItem(
                    i, j, QTableWidgetItem(str(elem)))

        f = True
        result = cur.execute("""SELECT title, id FROM sessions""").fetchall()
        self.combo_box_1.clear()
        for i in result:
            if f:
                self.chosen1 = f'{i[0]}id:{i[1]}'
                f = False
            self.combo_box_1.addItem(f'{i[0]}id:{i[1]}')

    def handle_activated_1(self, index):
        self.chosen1 = self.combo_box_1.itemText(index)

    def handle_activated_2(self, index):
        self.chosen2 = self.combo_box_2.itemText(index)


def except_hook(cls, exception, traceback):
    sys.__excepthook__(cls, exception, traceback)


if __name__ == '__main__':
    app = QApplication(sys.argv)
    start_window = StartWindow()
    start_window.show()
    sys.except_hook = except_hook
    sys.exit(app.exec_())
